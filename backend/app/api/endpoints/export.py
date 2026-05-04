from fastapi import APIRouter
from fastapi.responses import StreamingResponse
import io, json
import google.generativeai as genai
from fpdf import FPDF
from docx import Document
from pptx import Presentation
from app.schemas import ExportRequest

router = APIRouter()

# --- EXPORT ROUTES ---

@router.post("/export/word")
def export_word_ai(req: ExportRequest):
    model = genai.GenerativeModel("gemini-3-flash-preview")
    
    prompt = f"""
    You are an expert document formatter. Analyze the text below and organize it into a professional, well-structured document.
    Create a logical flow using a main title, section headings, paragraphs, and bullet points where appropriate.
    
    Return the result STRICTLY as a JSON object with this exact structure:
    {{
        "title": "Main Document Title",
        "sections": [
            {{
                "heading": "Section Heading Here",
                "paragraphs": ["Paragraph 1 text...", "Paragraph 2 text..."],
                "bullets": ["Optional bullet 1", "Optional bullet 2"]
            }}
        ]
    }}
    
    Content to process:
    {req.content}
    """
    
    try:
        response = model.generate_content(
            prompt,
            generation_config={"response_mime_type": "application/json"}
        )
        
        raw_text = response.text.strip()
        data = json.loads(raw_text)
        
    except Exception as e:
        print(f"Word Generation Error: {e}")
        # Fallback if AI fails
        data = {
            "title": req.title if req.title else "VidioMind Summary",
            "sections": [
                {"heading": "Summary", "paragraphs": [req.content], "bullets": []}
            ]
        }

    # --- WORD CREATION ---
    doc = Document()
    
    # Main Title
    doc.add_heading(data.get("title", req.title), 0)
    
    # Add Sections
    for sec in data.get("sections", []):
        if sec.get("heading"):
            doc.add_heading(str(sec["heading"]), level=1)
            
        for para in sec.get("paragraphs", []):
            if str(para).strip():
                doc.add_paragraph(str(para).strip())
                
        for bullet in sec.get("bullets", []):
            if str(bullet).strip():
                doc.add_paragraph(str(bullet).strip(), style='List Bullet')
                
    # Return File
    file_stream = io.BytesIO()
    doc.save(file_stream)
    file_stream.seek(0)
    
    return StreamingResponse(
        file_stream, 
        media_type="application/vnd.openxmlformats-officedocument.wordprocessingml.document", 
        headers={"Content-Disposition": f"attachment; filename=VidioMind_Document.docx"}
    )

@router.post("/export/pdf")
def export_pdf_ai(req: ExportRequest):
    model = genai.GenerativeModel("gemini-3-flash-preview")
    
    prompt = f"""
    You are an expert document formatter. Analyze the text below and organize it into a professional, well-structured document.
    Create a logical flow using a main title, section headings, paragraphs, and bullet points.
    
    Return the result STRICTLY as a JSON object with this exact structure:
    {{
        "title": "Generate a simple, direct, and factual title based on the text(half line only)",
        "sections": [
            {{
                "heading": "Section Heading Here",
                "paragraphs": ["Paragraph 1 text...", "Paragraph 2 text..."],
                "bullets": ["Optional bullet 1", "Optional bullet 2"]
            }}
        ]
    }}
    
    Content to process:
    {req.content}
    """
    
    try:
        response = model.generate_content(
            prompt,
            generation_config={"response_mime_type": "application/json"}
        )
        data = json.loads(response.text.strip())
    except Exception as e:
        print(f"PDF Generation Error: {e}")
        data = {
            "title": req.title if req.title else "VidioMind Summary",
            "sections": [{"heading": "Content", "paragraphs": [req.content], "bullets": []}]
        }

    # --- PDF CREATION ---
    pdf = FPDF()
    pdf.add_page()
    pdf.set_margins(15, 20, 15)
    
    # POWERFUL CLEANER: '?' wale issue ko rokne ke liye
    def clean_text(text):
        if not text:
            return ""
        text = str(text)
        text = text.replace('"', '"').replace('"', '"')
        text = text.replace("'", "'").replace("'", "'")
        text = text.replace('\u2013', '-').replace('\u2014', '-')
        # Unicode bullet hata rahe hain taake FPDF error na de
        text = text.replace('\u2022', '') 
        return text.encode('latin-1', 'replace').decode('latin-1')

    # Main Title
    pdf.set_font("helvetica", style="B", size=18)
    # Yahan multi_cell lagaya hai taake title jitna bhi lamba ho, wo proper center aligned lines mein tut jaye
    pdf.multi_cell(0, 10, text=clean_text(data.get("title", req.title)), align="C")
    pdf.ln(5)
    
    # Add Sections
    for sec in data.get("sections", []):
        if sec.get("heading"):
            pdf.set_font("helvetica", style="B", size=14)
            pdf.set_text_color(37, 99, 235) 
            pdf.multi_cell(0, 8, text=clean_text(sec["heading"]))
            pdf.ln(2)
            
        pdf.set_font("helvetica", size=11)
        pdf.set_text_color(71, 85, 105) 
        
        for para in sec.get("paragraphs", []):
            if str(para).strip():
                pdf.multi_cell(0, 7, text=clean_text(para))
                pdf.ln(2)
                
        for bullet in sec.get("bullets", []):
            if str(bullet).strip():
                # FIX: Yahan chr(149) use kiya hy jo proper '•' gol bullet print karega
                pdf.multi_cell(0, 7, text=clean_text(f"{chr(149)} {bullet}"))
                pdf.ln(1)
                
        pdf.ln(4) 
        
    pdf_bytes = bytes(pdf.output())
    return StreamingResponse(
        io.BytesIO(pdf_bytes), 
        media_type="application/pdf", 
        headers={"Content-Disposition": f"attachment; filename=VidioMind_Document.pdf"}
    )

@router.post("/export/ppt")
def export_ppt_ai(req: ExportRequest):
    model = genai.GenerativeModel("gemini-3-flash-preview")
    
    prompt = f"""
    You are an expert presentation maker. Analyze the text below and convert ALL of it into a presentation.
    Create AS MANY slides as needed (e.g., 5, 10, or more) to cover the entire content comprehensively.
    Break down large paragraphs into concise bullet points (maximum 5 bullets per slide).
    
    Return the result STRICTLY as a JSON object with a single key "slides" containing an array of slide objects.
    Format Example:
    {{
        "slides": [
            {{"title": "First Topic", "bullets": ["Point 1", "Point 2"]}},
            {{"title": "Second Topic", "bullets": ["Point 1", "Point 2", "Point 3"]}}
        ]
    }}
    
    Content to process:
    {req.content}
    """
    
    try:
        # Force Gemini to return strictly valid JSON
        response = model.generate_content(
            prompt,
            generation_config={"response_mime_type": "application/json"}
        )
        
        raw_text = response.text.strip()
        
        # Extra safety parser just in case
        if "```json" in raw_text:
            raw_text = raw_text.split("```json")[1].split("```")[0].strip()
            
        start_idx = raw_text.find('{')
        end_idx = raw_text.rfind('}') + 1
        if start_idx != -1 and end_idx != -1:
            raw_text = raw_text[start_idx:end_idx]
            
        data = json.loads(raw_text)
        
        # Handle the slides array safely
        if isinstance(data, dict) and "slides" in data:
            slides_data = data["slides"]
        elif isinstance(data, list):
            slides_data = data
        else:
            raise ValueError("Unexpected JSON format received")
            
    except Exception as e:
        print(f"PPT Generation Error: {e}")
        # Agar bilkul hi masla aa jaye tab ye slide chalegi
        slides_data = [
            {"title": "Error in AI Generation", "bullets": ["Sorry, the AI could not format this text properly.", str(e)]},
            {"title": "Raw Text", "bullets": [req.content[:200] + "..."]}
        ]

    # --- PPT CREATION ---
    prs = Presentation()
    
    # Title Slide
    title_slide_layout = prs.slide_layouts[0]
    slide = prs.slides.add_slide(title_slide_layout)
    slide.shapes.title.text = req.title if req.title else "VidioMind Presentation"
    slide.placeholders[1].text = "Generated by VidioMind AI"
    
    # Bullet Slides (All generated content)
    bullet_slide_layout = prs.slide_layouts[1]
    
    for slide_info in slides_data:
        slide = prs.slides.add_slide(bullet_slide_layout)
        shapes = slide.shapes
        
        # Title
        if shapes.title:
            shapes.title.text = str(slide_info.get("title", "Slide Title"))
        
        # Bullets
        body_shape = shapes.placeholders[1]
        tf = body_shape.text_frame
        tf.clear() # Clear default prompt text
        
        bullets = slide_info.get("bullets", [])
        for i, bullet_text in enumerate(bullets):
            p = tf.add_paragraph() if i > 0 else tf.paragraphs[0]
            p.text = str(bullet_text).strip()
            p.level = 0 # Proper bullet level
                
    # Return File
    file_stream = io.BytesIO()
    prs.save(file_stream)
    file_stream.seek(0)
    
    return StreamingResponse(
        file_stream, 
        media_type="application/vnd.openxmlformats-officedocument.presentationml.presentation", 
        headers={"Content-Disposition": f"attachment; filename=VidioMind_Presentation.pptx"}
    )