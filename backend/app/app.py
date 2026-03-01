import json
import os
import time
import yt_dlp
from google import genai
import google.generativeai as genai
from fastapi import FastAPI, Depends, BackgroundTasks, HTTPException, status
from fastapi.middleware.cors import CORSMiddleware
from sqlalchemy.orm import Session
from pydantic import BaseModel
from passlib.context import CryptContext
from dotenv import load_dotenv
import models, database
import io
from fastapi.responses import StreamingResponse
from fpdf import FPDF
from docx import Document
from pptx import Presentation
import json
from pptx import Presentation
from pptx.util import Inches

# Export ke liye naya schema add karein
class ExportRequest(BaseModel):
    title: str
    content: str

# --- CONFIGURATION ---
load_dotenv()
GEMINI_KEY = os.getenv("GEMINI_API_KEY")

if not GEMINI_KEY:
    print("Error: GEMINI_API_KEY not found in .env file")

genai.configure(api_key=GEMINI_KEY)

app = FastAPI()

# Database Setup
models.Base.metadata.create_all(bind=database.engine)

# CORS Setup
app.add_middleware(
    CORSMiddleware,
    allow_origins=["*"],
    allow_credentials=True,
    allow_methods=["*"],
    allow_headers=["*"],
)

# Password Hashing
pwd_context = CryptContext(schemes=["bcrypt"], deprecated="auto")

def verify_password(plain, hashed): return pwd_context.verify(plain, hashed)
def get_password_hash(password): return pwd_context.hash(password)

# --- SCHEMAS ---
class UserBase(BaseModel):
    email: str
    password: str
    username: str = "User"

class UserResponse(BaseModel):
    user_id: int
    email: str
    username: str

# --- AUTH ROUTES ---
@app.post("/api/register", response_model=UserResponse)
def register(user: UserBase, db: Session = Depends(database.get_db)):
    if db.query(models.User).filter(models.User.email == user.email).first():
        raise HTTPException(status_code=400, detail="Email already registered")
    
    new_user = models.User(
        email=user.email, 
        password_hash=get_password_hash(user.password), 
        username=user.username
    )
    db.add(new_user)
    db.commit()
    db.refresh(new_user)
    return new_user

@app.post("/api/login", response_model=UserResponse)
def login(req: UserBase, db: Session = Depends(database.get_db)):
    user = db.query(models.User).filter(models.User.email == req.email).first()
    if not user or not verify_password(req.password, user.password_hash):
        raise HTTPException(status_code=401, detail="Invalid credentials")
    return user

# --- CORE FEATURES ---

# 1. SUMMARIZE VIDEO
@app.post("/api/summarize")
def submit_video(url: str, user_id: int, background_tasks: BackgroundTasks, db: Session = Depends(database.get_db)):
    # Thumbnail Extraction
    vid_id = "default"
    if "v=" in url:
        try: vid_id = url.split("v=")[1].split("&")[0]
        except: pass
    elif "youtu.be" in url:
        try: vid_id = url.split("/")[-1]
        except: pass
    
    thumbnail = f"https://img.youtube.com/vi/{vid_id}/mqdefault.jpg"

    # Create Project
    project = models.VideoProject(
        user_id=user_id, source_url=url, title="Processing Video...",
        thumbnail_url=thumbnail, status="Queued"
    )
    db.add(project)
    db.commit()
    db.refresh(project)

    # Tracker
    tracker = models.ProcessingRequest(project_id=project.project_id, job_status="Queued", progress_percent=0)
    db.add(tracker)
    db.commit()

    # Start Processing
    background_tasks.add_task(process_video_with_gemini, project.project_id, url, db)
    
    return {"project_id": project.project_id}

# 2. VIEW HISTORY
@app.get("/api/history/{user_id}")
def get_history(user_id: int, db: Session = Depends(database.get_db)):
    return db.query(models.VideoProject).filter(models.VideoProject.user_id == user_id).order_by(models.VideoProject.created_at.desc()).all()

# 3. GET RESULT
@app.get("/api/status/{project_id}")
def get_status(project_id: int, db: Session = Depends(database.get_db)):
    tracker = db.query(models.ProcessingRequest).filter_by(project_id=project_id).first()
    if not tracker: raise HTTPException(status_code=404)
    
    res = {
        "status": tracker.job_status, 
        "progress": tracker.progress_percent, 
        "summary": None,
        "blog": None
    }
    
    if tracker.job_status == "Completed":
        summ = db.query(models.Summary).filter_by(project_id=project_id).first()
        if summ: 
            res["summary"] = summ.concise_summary
            res["blog"] = summ.blog_post_content
            
    return res

# --- EXPORT ROUTES ---

@app.post("/api/export/word")
def export_word_ai(req: ExportRequest):
    model = genai.GenerativeModel("gemini-3-flash-preview")
    
    prompt = f"""
    You are an expert document formatter. Analyze the text below and organize it into a professional, well-structured document.
    Create a logical flow using a main title, section headings, paragraphs, and bullet points where appropriate.
    
    Return the result STRICTLY as a JSON object with this exact structure:
    {{
        "title": "Main Document Title",
        "sections": [
            {{
                "heading": "Section Heading Here",
                "paragraphs": ["Paragraph 1 text...", "Paragraph 2 text..."],
                "bullets": ["Optional bullet 1", "Optional bullet 2"]
            }}
        ]
    }}
    
    Content to process:
    {req.content}
    """
    
    try:
        response = model.generate_content(
            prompt,
            generation_config={"response_mime_type": "application/json"}
        )
        
        raw_text = response.text.strip()
        data = json.loads(raw_text)
        
    except Exception as e:
        print(f"Word Generation Error: {e}")
        # Fallback if AI fails
        data = {
            "title": req.title if req.title else "VidioMind Summary",
            "sections": [
                {"heading": "Summary", "paragraphs": [req.content], "bullets": []}
            ]
        }

    # --- WORD CREATION ---
    doc = Document()
    
    # Main Title
    doc.add_heading(data.get("title", req.title), 0)
    
    # Add Sections
    for sec in data.get("sections", []):
        if sec.get("heading"):
            doc.add_heading(str(sec["heading"]), level=1)
            
        for para in sec.get("paragraphs", []):
            if str(para).strip():
                doc.add_paragraph(str(para).strip())
                
        for bullet in sec.get("bullets", []):
            if str(bullet).strip():
                doc.add_paragraph(str(bullet).strip(), style='List Bullet')
                
    # Return File
    file_stream = io.BytesIO()
    doc.save(file_stream)
    file_stream.seek(0)
    
    return StreamingResponse(
        file_stream, 
        media_type="application/vnd.openxmlformats-officedocument.wordprocessingml.document", 
        headers={"Content-Disposition": f"attachment; filename=VidioMind_Document.docx"}
    )

@app.post("/api/export/pdf")
def export_pdf_ai(req: ExportRequest):
    model = genai.GenerativeModel("gemini-3-flash-preview")
    
    prompt = f"""
    You are an expert document formatter. Analyze the text below and organize it into a professional, well-structured document.
    Create a logical flow using a main title, section headings, paragraphs, and bullet points.
    
    Return the result STRICTLY as a JSON object with this exact structure:
    {{
        "title": "Generate a simple, direct, and factual title based on the text(half line only)",
        "sections": [
            {{
                "heading": "Section Heading Here",
                "paragraphs": ["Paragraph 1 text...", "Paragraph 2 text..."],
                "bullets": ["Optional bullet 1", "Optional bullet 2"]
            }}
        ]
    }}
    
    Content to process:
    {req.content}
    """
    
    try:
        response = model.generate_content(
            prompt,
            generation_config={"response_mime_type": "application/json"}
        )
        data = json.loads(response.text.strip())
    except Exception as e:
        print(f"PDF Generation Error: {e}")
        data = {
            "title": req.title if req.title else "VidioMind Summary",
            "sections": [{"heading": "Content", "paragraphs": [req.content], "bullets": []}]
        }

    # --- PDF CREATION ---
    pdf = FPDF()
    pdf.add_page()
    pdf.set_margins(15, 20, 15)
    
    # POWERFUL CLEANER: '?' wale issue ko rokne ke liye
    def clean_text(text):
        if not text:
            return ""
        text = str(text)
        text = text.replace('"', '"').replace('"', '"')
        text = text.replace("'", "'").replace("'", "'")
        text = text.replace('\u2013', '-').replace('\u2014', '-')
        # Unicode bullet hata rahe hain taake FPDF error na de
        text = text.replace('\u2022', '') 
        return text.encode('latin-1', 'replace').decode('latin-1')

    # Main Title
    pdf.set_font("helvetica", style="B", size=18)
    # Yahan multi_cell lagaya hai taake title jitna bhi lamba ho, wo proper center aligned lines mein tut jaye
    pdf.multi_cell(0, 10, text=clean_text(data.get("title", req.title)), align="C")
    pdf.ln(5)
    
    # Add Sections
    for sec in data.get("sections", []):
        if sec.get("heading"):
            pdf.set_font("helvetica", style="B", size=14)
            pdf.set_text_color(37, 99, 235) 
            pdf.multi_cell(0, 8, text=clean_text(sec["heading"]))
            pdf.ln(2)
            
        pdf.set_font("helvetica", size=11)
        pdf.set_text_color(71, 85, 105) 
        
        for para in sec.get("paragraphs", []):
            if str(para).strip():
                pdf.multi_cell(0, 7, text=clean_text(para))
                pdf.ln(2)
                
        for bullet in sec.get("bullets", []):
            if str(bullet).strip():
                # FIX: Yahan chr(149) use kiya hy jo proper '•' gol bullet print karega
                pdf.multi_cell(0, 7, text=clean_text(f"{chr(149)} {bullet}"))
                pdf.ln(1)
                
        pdf.ln(4) 
        
    pdf_bytes = bytes(pdf.output())
    return StreamingResponse(
        io.BytesIO(pdf_bytes), 
        media_type="application/pdf", 
        headers={"Content-Disposition": f"attachment; filename=VidioMind_Document.pdf"}
    )

@app.post("/api/export/ppt")
def export_ppt_ai(req: ExportRequest):
    model = genai.GenerativeModel("gemini-3-flash-preview")
    
    prompt = f"""
    You are an expert presentation maker. Analyze the text below and convert ALL of it into a presentation.
    Create AS MANY slides as needed (e.g., 5, 10, or more) to cover the entire content comprehensively.
    Break down large paragraphs into concise bullet points (maximum 5 bullets per slide).
    
    Return the result STRICTLY as a JSON object with a single key "slides" containing an array of slide objects.
    Format Example:
    {{
        "slides": [
            {{"title": "First Topic", "bullets": ["Point 1", "Point 2"]}},
            {{"title": "Second Topic", "bullets": ["Point 1", "Point 2", "Point 3"]}}
        ]
    }}
    
    Content to process:
    {req.content}
    """
    
    try:
        # Force Gemini to return strictly valid JSON
        response = model.generate_content(
            prompt,
            generation_config={"response_mime_type": "application/json"}
        )
        
        raw_text = response.text.strip()
        
        # Extra safety parser just in case
        if "```json" in raw_text:
            raw_text = raw_text.split("```json")[1].split("```")[0].strip()
            
        start_idx = raw_text.find('{')
        end_idx = raw_text.rfind('}') + 1
        if start_idx != -1 and end_idx != -1:
            raw_text = raw_text[start_idx:end_idx]
            
        data = json.loads(raw_text)
        
        # Handle the slides array safely
        if isinstance(data, dict) and "slides" in data:
            slides_data = data["slides"]
        elif isinstance(data, list):
            slides_data = data
        else:
            raise ValueError("Unexpected JSON format received")
            
    except Exception as e:
        print(f"PPT Generation Error: {e}")
        # Agar bilkul hi masla aa jaye tab ye slide chalegi
        slides_data = [
            {"title": "Error in AI Generation", "bullets": ["Sorry, the AI could not format this text properly.", str(e)]},
            {"title": "Raw Text", "bullets": [req.content[:200] + "..."]}
        ]

    # --- PPT CREATION ---
    prs = Presentation()
    
    # Title Slide
    title_slide_layout = prs.slide_layouts[0]
    slide = prs.slides.add_slide(title_slide_layout)
    slide.shapes.title.text = req.title if req.title else "VidioMind Presentation"
    slide.placeholders[1].text = "Generated by VidioMind AI"
    
    # Bullet Slides (All generated content)
    bullet_slide_layout = prs.slide_layouts[1]
    
    for slide_info in slides_data:
        slide = prs.slides.add_slide(bullet_slide_layout)
        shapes = slide.shapes
        
        # Title
        if shapes.title:
            shapes.title.text = str(slide_info.get("title", "Slide Title"))
        
        # Bullets
        body_shape = shapes.placeholders[1]
        tf = body_shape.text_frame
        tf.clear() # Clear default prompt text
        
        bullets = slide_info.get("bullets", [])
        for i, bullet_text in enumerate(bullets):
            p = tf.add_paragraph() if i > 0 else tf.paragraphs[0]
            p.text = str(bullet_text).strip()
            p.level = 0 # Proper bullet level
                
    # Return File
    file_stream = io.BytesIO()
    prs.save(file_stream)
    file_stream.seek(0)
    
    return StreamingResponse(
        file_stream, 
        media_type="application/vnd.openxmlformats-officedocument.presentationml.presentation", 
        headers={"Content-Disposition": f"attachment; filename=VidioMind_Presentation.pptx"}
    )

# --- GEMINI AI ENGINE (Advanced Formatting Fix) ---
def process_video_with_gemini(pid: int, url: str, db: Session):
    audio_file = None
    try:
        # Step 1: Download
        update_status(pid, "Downloading Audio...", 10, db)
        
        ydl_opts = {
            'format': 'bestaudio[ext=m4a]/bestaudio/best',
            'outtmpl': f'temp_{pid}.%(ext)s',
            'quiet': True,
            'nocheckcertificate': True,
            'socket_timeout': 60,
            'retries': 10,
        }
        
        with yt_dlp.YoutubeDL(ydl_opts) as ydl:
            info = ydl.extract_info(url, download=True)
            video_title = info.get('title', 'Unknown Video')
            thumbnail = info.get('thumbnail', '')
            ext = info.get('ext', 'm4a')
            audio_file = f"temp_{pid}.{ext}"

        # Update Meta
        proj = db.query(models.VideoProject).filter_by(project_id=pid).first()
        proj.title = video_title
        proj.thumbnail_url = thumbnail
        db.commit()

        # Step 2: Upload
        update_status(pid, "Uploading to Gemini AI...", 40, db)
        myfile = genai.upload_file(audio_file)
        
        while myfile.state.name == "PROCESSING":
            time.sleep(4)
            myfile = genai.get_file(myfile.name)
            
        if myfile.state.name == "FAILED":
            raise Exception("Gemini failed to process the audio file.")

        # Step 3: Generate
        update_status(pid, "Generating Clean Content...", 75, db)
        
        model = genai.GenerativeModel("gemini-3-flash-preview")
        
        prompt = """
        Analyze this audio. Return a strictly valid JSON object.
        1. "summary": A simple string with bullet points (use • symbol).
        2. "blog_post": A simple string containing a Title and Paragraphs separated by newlines.
        
        Important: Do NOT use nested objects for blog_post (like {"title": "..."}). Just return plain text strings for both fields.
        Return ONLY raw JSON.
        """
        
        result = model.generate_content([myfile, prompt])
        
        # --- ADVANCED CLEANER ---
        final_summary = ""
        final_blog = ""
        
        try:
            raw_text = result.text.strip()
            # Markdown cleaning
            if "```json" in raw_text:
                raw_text = raw_text.replace("```json", "").replace("```", "")
            
            # Braces extraction
            start = raw_text.find('{')
            end = raw_text.rfind('}') + 1
            if start != -1 and end != -1:
                raw_text = raw_text[start:end]

            data = json.loads(raw_text)
            
            # --- FIX: Handle Nested Structures (Summary) ---
            s_data = data.get("summary", "No summary.")
            if isinstance(s_data, list):
                final_summary = "\n".join([f"• {str(item)}" for item in s_data])
            elif isinstance(s_data, dict):
                # Agar ghalati se dict aa jaye
                final_summary = "\n".join([f"• {v}" for v in s_data.values()])
            else:
                final_summary = str(s_data)

            # --- FIX: Handle Nested Structures (Blog) ---
            # Screenshot mein {title: "", paragraphs: []} tha, usay handle karte hain
            b_data = data.get("blog_post", "No blog.")
            
            if isinstance(b_data, dict):
                # Title aur Paragraphs ko merge karo
                title = b_data.get("title", "Blog Post")
                content = b_data.get("paragraphs", [])
                if isinstance(content, list):
                    content = "\n\n".join([str(p) for p in content])
                elif isinstance(content, str):
                    # Kabhi kabhi content string hota hai lekin 'content' key mein
                    content = b_data.get("content", content)
                
                final_blog = f"{title}\n\n{content}"
            
            elif isinstance(b_data, list):
                final_blog = "\n\n".join([str(item) for item in b_data])
            else:
                final_blog = str(b_data)

            # Final Cleanup
            final_summary = final_summary.strip()
            final_blog = final_blog.strip()

        except json.JSONDecodeError:
            final_summary = result.text
            final_blog = "Format Error: Could not parse blog."

        # Step 4: Save
        summ = models.Summary(
            project_id=pid, 
            concise_summary=final_summary,
            blog_post_content=final_blog
        )
        db.add(summ)
        
        update_status(pid, "Completed", 100, db)
        proj.status = "Completed"
        db.commit()

        # Cleanup
        if audio_file and os.path.exists(audio_file): os.remove(audio_file)

    except Exception as e:
        print(f"Error details: {e}")
        db.rollback()
        
        try:
            update_status(pid, "Failed", 0, db)
            proj = db.query(models.VideoProject).filter_by(project_id=pid).first()
            if proj:
                proj.status = "Failed"
                db.commit()
        except: pass

        if audio_file and os.path.exists(audio_file): os.remove(audio_file)

def update_status(pid, status, percent, db):
    try:
        req = db.query(models.ProcessingRequest).filter_by(project_id=pid).first()
        if req:
            req.job_status = status
            req.progress_percent = percent
            db.commit()
    except:
        db.rollback()